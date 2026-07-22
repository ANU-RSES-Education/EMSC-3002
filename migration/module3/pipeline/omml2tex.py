#!/usr/bin/env python3
"""Convert OMML (Office Math, <m:oMath>) lxml elements to LaTeX.
Handles the constructs used in the Module-3 stress/strain/rheology decks:
fractions, sub/superscripts, radicals, delimiters, n-ary (sum/int), matrices,
accents (vector/hat/dot/bar), functions. Greek/unicode symbols pass through
(MathJax renders them in math mode)."""
import re, unicodedata
from lxml import etree

M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
def q(tag): return f"{{{M}}}{tag}"

# Greek unicode -> LaTeX (after NFKC normalisation of math-italic code points).
GREEK = {
    "α": r"\alpha ", "β": r"\beta ", "γ": r"\gamma ", "δ": r"\delta ",
    "ε": r"\epsilon ", "ϵ": r"\epsilon ", "ζ": r"\zeta ", "η": r"\eta ",
    "θ": r"\theta ", "ϑ": r"\vartheta ", "ι": r"\iota ", "κ": r"\kappa ",
    "λ": r"\lambda ", "μ": r"\mu ", "ν": r"\nu ", "ξ": r"\xi ",
    "π": r"\pi ", "ϖ": r"\varpi ", "ρ": r"\rho ", "ϱ": r"\varrho ",
    "σ": r"\sigma ", "ς": r"\varsigma ", "τ": r"\tau ", "υ": r"\upsilon ",
    "φ": r"\phi ", "ϕ": r"\varphi ", "χ": r"\chi ", "ψ": r"\psi ",
    "ω": r"\omega ",
    "Γ": r"\Gamma ", "Δ": r"\Delta ", "Θ": r"\Theta ", "Λ": r"\Lambda ",
    "Ξ": r"\Xi ", "Π": r"\Pi ", "Σ": r"\Sigma ", "Υ": r"\Upsilon ",
    "Φ": r"\Phi ", "Ψ": r"\Psi ", "Ω": r"\Omega ",
}
def norm_text(s):
    """NFKC folds math-italic/bold code points to plain letters; then Greek->LaTeX."""
    s = unicodedata.normalize("NFKC", s or "")
    s = s.replace(" ", " ")            # non-breaking space padding
    s = "".join(c for c in s if c not in INVISIBLE)
    s = s.replace("∆", r"\Delta ")  # ∆ increment sign -> Delta
    return "".join(GREEK.get(c, c) for c in s)

INVISIBLE = {"⁡", "⁢", "⁣", "⁤"}  # function-application / invisible-times etc.

# accent char -> latex wrapper
ACC = {"⃗": "vec", "̂": "hat", "^": "hat", "̇": "dot", "̈": "ddot",
       "̃": "tilde", "~": "tilde", "̄": "bar", "¯": "bar", "ˉ": "bar",
       "→": "vec"}
# a few operators worth normalising; most unicode passes straight through to MathJax
OPS = {"×": r" \times ", "·": r" \cdot ", "−": "-", "≈": r" \approx ",
       "≤": r" \le ", "≥": r" \ge ", "≠": r" \ne ", "±": r" \pm ",
       "≅": r" \cong ", "∝": r" \propto ", "∂": r" \partial ",
       "∇": r" \nabla ", "∑": r" \sum ", "∫": r" \int "}

def _text(el):
    return norm_text(el.text or "")

def conv(el):
    """Recursively convert an OMML element's children to a LaTeX string."""
    out = []
    for ch in el:
        tag = etree.QName(ch).localname
        out.append(HANDLERS.get(tag, _passthrough)(ch))
    return "".join(out)

def _passthrough(el):
    return conv(el)

def h_r(el):  # run: gather m:t text
    s = "".join(norm_text(t.text or "") for t in el.iter(q("t")))
    for k, v in OPS.items():
        s = s.replace(k, v)
    return s

def _child(el, name):
    c = el.find(q(name))
    return conv(c) if c is not None else ""

def h_f(el):  # fraction
    num, den = _child(el, "num"), _child(el, "den")
    pr = el.find(q("fPr"))
    typ = None
    if pr is not None:
        t = pr.find(q("type"))
        if t is not None: typ = t.get(q("val"))
    if typ == "noBar":
        return "{" + num + r" \atop " + den + "}"
    if typ == "lin":
        return f"{num}/{den}"
    return r"\frac{" + num + "}{" + den + "}"

def h_ssup(el):
    return "{" + _child(el, "e") + "}^{" + _child(el, "sup") + "}"
def h_ssub(el):
    return "{" + _child(el, "e") + "}_{" + _child(el, "sub") + "}"
def h_ssubsup(el):
    return "{" + _child(el, "e") + "}_{" + _child(el, "sub") + "}^{" + _child(el, "sup") + "}"
def h_spre(el):
    return "{}_{" + _child(el, "sub") + "}^{" + _child(el, "sup") + "}" + _child(el, "e")

def h_rad(el):
    deg = _child(el, "deg")
    pr = el.find(q("radPr"))
    hide = pr is not None and (pr.find(q("degHide")) is not None and pr.find(q("degHide")).get(q("val")) in ("1", "true", None))
    e = _child(el, "e")
    if deg and not hide:
        return r"\sqrt[" + deg + "]{" + e + "}"
    return r"\sqrt{" + e + "}"

DELIM = {"[": ("[", "]"), "{": (r"\{", r"\}"), "(": ("(", ")"), "|": ("|", "|"),
         "‖": (r"\|", r"\|"), "⟨": (r"\langle", r"\rangle"), "": (".", ".")}
def h_d(el):  # delimiter
    pr = el.find(q("dPr"))
    beg, end, sep = "(", ")", "|"
    if pr is not None:
        b = pr.find(q("begChr")); e = pr.find(q("endChr")); s = pr.find(q("sepChr"))
        if b is not None: beg = b.get(q("val")) or ""
        if e is not None: end = e.get(q("val")) or ""
        if s is not None: sep = s.get(q("val")) or "|"
    lb = {"[": "[", "{": r"\{", "(": "(", "|": "|", "": ".", "⌈": r"\lceil",
          "⌊": r"\lfloor", "⟨": r"\langle"}.get(beg, beg or ".")
    rb = {"]": "]", "}": r"\}", ")": ")", "|": "|", "": ".", "⌉": r"\rceil",
          "⌋": r"\rfloor", "⟩": r"\rangle"}.get(end, end or ".")
    es = [conv(e) for e in el.findall(q("e"))]
    body = (" " + sep + " ").join(es)
    return r"\left" + lb + " " + body + r" \right" + rb

def h_nary(el):  # sum / integral / product
    pr = el.find(q("naryPr"))
    chr_ = "∫"
    subhid = suphid = False
    if pr is not None:
        c = pr.find(q("chr"))
        if c is not None: chr_ = c.get(q("val")) or chr_
        sh = pr.find(q("subHide")); ph = pr.find(q("supHide"))
        subhid = sh is not None and sh.get(q("val")) in ("1", "true")
        suphid = ph is not None and ph.get(q("val")) in ("1", "true")
    op = {"∑": r"\sum", "∫": r"\int", "∏": r"\prod",
          "∬": r"\iint", "∭": r"\iiint", "∐": r"\coprod"}.get(chr_, r"\int")
    sub = _child(el, "sub"); sup = _child(el, "sup"); e = _child(el, "e")
    s = op
    if not subhid and sub: s += "_{" + sub + "}"
    if not suphid and sup: s += "^{" + sup + "}"
    return s + " " + e

def h_acc(el):
    pr = el.find(q("accPr"))
    ch = "̂"
    if pr is not None:
        c = pr.find(q("chr"))
        if c is not None: ch = c.get(q("val")) or ch
    wrap = ACC.get(ch, "hat")
    return "\\" + wrap + "{" + _child(el, "e") + "}"

def h_bar(el):
    pr = el.find(q("barPr"))
    pos = "top"
    if pr is not None:
        p = pr.find(q("pos"))
        if p is not None: pos = p.get(q("val")) or "top"
    return ("\\underline{" if pos == "bot" else "\\overline{") + _child(el, "e") + "}"

def h_func(el):
    name = _child(el, "fName").strip()
    known = {"sin","cos","tan","cot","sec","csc","arcsin","arccos","arctan",
             "sinh","cosh","tanh","log","ln","lim","exp","min","max","det","gcd"}
    fn = "\\" + name if name in known else r"\operatorname{" + name + "}"
    return fn + "(" + _child(el, "e") + ")" if False else fn + " " + _child(el, "e")

def h_limlow(el):
    return "{" + _child(el, "e") + "}_{" + _child(el, "lim") + "}"
def h_limupp(el):
    return "{" + _child(el, "e") + "}^{" + _child(el, "lim") + "}"

def h_m(el):  # matrix
    rows = []
    for mr in el.findall(q("mr")):
        cells = [conv(e) for e in mr.findall(q("e"))]
        rows.append(" & ".join(cells))
    # \cr (not \\) for row breaks: reveal's markdown pass mangles \\ before KaTeX runs
    body = " \\cr ".join(rows)
    return r"\begin{matrix} " + body + r" \end{matrix}"

def h_groupChr(el):
    return _child(el, "e")

HANDLERS = {
    "r": h_r, "t": lambda e: _text(e), "f": h_f, "sSup": h_ssup, "sSub": h_ssub,
    "sSubSup": h_ssubsup, "sPre": h_spre, "rad": h_rad, "d": h_d, "nary": h_nary,
    "acc": h_acc, "bar": h_bar, "func": h_func, "limLow": h_limlow, "limUpp": h_limupp,
    "m": h_m, "e": _passthrough, "num": _passthrough, "den": _passthrough,
    "sup": _passthrough, "sub": _passthrough, "deg": _passthrough, "lim": _passthrough,
    "fName": _passthrough, "groupChr": h_groupChr, "box": _passthrough,
    "borderBox": _passthrough, "eqArr": _passthrough,
}

_TRIG = re.compile(r"(?<![\\A-Za-z])(sinh|cosh|tanh|arcsin|arccos|arctan|"
                   r"sin|cos|tan|cot|sec|csc|log|ln|exp)(?![A-Za-z])")
def cleanup(s):
    s = re.sub(r"\s+", " ", s).strip()
    s = s.replace("{ ", "{").replace(" }", "}")
    s = _TRIG.sub(lambda m: "\\" + m.group(1) + " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def omml_to_latex(omath_el):
    """omath_el: an <m:oMath> lxml element. Returns a LaTeX string (no $ delimiters)."""
    return cleanup(conv(omath_el))

if __name__ == "__main__":
    import sys
    from pptx import Presentation
    # quick test harness: dump LaTeX for oMath in given slide numbers of a pptx
    pptx, slidenos = sys.argv[1], [int(x) for x in sys.argv[2:]]
    prs = Presentation(pptx)
    for i, slide in enumerate(prs.slides, 1):
        if slidenos and i not in slidenos: continue
        found = []
        for sh in slide.shapes:
            if not sh.has_text_frame: continue
            for om in sh.text_frame._txBody.iter(q("oMath")):
                found.append(omml_to_latex(om))
        if found:
            print(f"=== slide {i} ===")
            for f in found:
                print("  $" + f + "$")
