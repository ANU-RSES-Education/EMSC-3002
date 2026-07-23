#!/usr/bin/env python3
"""XML-level pptx extractor that recovers OMML equations.

The equation-bearing shapes are wrapped in <mc:AlternateContent><mc:Choice> (the
rich math version) + <mc:Fallback> (garbled text). python-pptx's slide.shapes
skips the Choice and exposes only the Fallback, which is why the first-pass
conversion dropped/scrambled every equation. Here we walk the slide XML directly,
take the Choice subtree, convert <m:oMath> to LaTeX, and interleave it with the
normal text runs (plus text super/subscripts -> <sup>/<sub>).

Emits the same manifest schema as extract_pptx.py so convert_deck.py is unchanged.
Usage: extract_pptx2.py <pptx> <out_dir>"""
import sys, json, re, pathlib, io
from pptx import Presentation
from PIL import Image
sys.path.insert(0, str(pathlib.Path(__file__).parent))
from omml2tex import omml_to_latex

MAXPX = 1600          # cap long side (slides never need more)
JPEG_Q = 85

def optimize_image(blob, ext):
    """Return (new_ext, bytes): downscale oversized images, convert TIFF/BMP to a
    web format, recompress. Keeps PNG (line art / transparency) as PNG; everything
    else becomes JPEG. Falls back to the original blob on any failure."""
    ext = ext.lower().lstrip(".")
    if ext in ("gif", "svg"):            # leave animations / vectors alone
        return ext, blob
    try:
        im = Image.open(io.BytesIO(blob))
        im.load()
    except Exception:
        return ext, blob
    # alpha only counts if it is actually used (many pptx PNGs are opaque RGBA)
    has_alpha = False
    if im.mode in ("RGBA", "LA"):
        has_alpha = im.getchannel("A").getextrema()[0] < 255
    elif im.mode == "P" and "transparency" in im.info:
        has_alpha = True
    w, h = im.size
    scale = min(1.0, MAXPX / max(w, h))
    if scale < 1.0:
        im = im.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.LANCZOS)

    if has_alpha:                       # transparency must stay PNG
        if im.mode not in ("RGBA", "LA", "P"):
            im = im.convert("RGBA")
        buf = io.BytesIO(); im.save(buf, format="PNG", optimize=True)
        cand = [("png", buf.getvalue())]
    else:                               # opaque: pick the smaller of JPEG / PNG
        rgb = im.convert("RGB") if im.mode != "RGB" else im
        jb = io.BytesIO(); rgb.save(jb, format="JPEG", quality=JPEG_Q, optimize=True, progressive=True)
        pb = io.BytesIO(); (im if im.mode in ("P", "L") else rgb).save(pb, format="PNG", optimize=True)
        cand = [("jpg", jb.getvalue()), ("png", pb.getvalue())]
    new_ext, out = min(cand, key=lambda c: len(c[1]))
    # if the source was already a web format and we didn't shrink it, keep the original
    if ext in ("png", "jpg", "jpeg") and scale == 1.0 and len(out) >= len(blob):
        return ext, blob
    return new_ext, out

P  = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
A  = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
MC = "{http://schemas.openxmlformats.org/markup-compatibility/2006}"
M  = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"

ATTR = re.compile(r"^\s*(source|image|photo|credit|courtesy|after|from|adapted|modified|©|copyright)\b", re.I)
URLRE = re.compile(r"(https?://|www\.|\.com|\.org|\.edu|\.au\b)", re.I)
AUTHYEAR = re.compile(r"^[A-Z][A-Za-z.\-]+(?:\s+(?:et al\.?|and|&|[A-Z][A-Za-z.\-]+))*,?\s*\(?(19|20)\d\d")

def local(el): return el.tag.split("}")[-1] if isinstance(el.tag, str) else ""
def pct(v, total): return round(100.0 * v / total, 1) if total else 0.0

# ---------- math + text run helpers ----------
def emit_math(node, restrict_choice=True):
    """Return list of LaTeX strings for oMath found under node.
    For mc:AlternateContent restrict to the Choice (skip the Fallback duplicate)."""
    if local(node) == "AlternateContent" and restrict_choice:
        choice = node.find(MC + "Choice")
        node = choice if choice is not None else node
    return [omml_to_latex(om) for om in node.iter(M + "oMath")]

def wrap_math(latex):
    if "\\begin{matrix}" in latex or len(latex) > 46:
        return None, f"$${latex}$$"          # display, own paragraph
    return f"${latex}$", None                # inline

def run_text(r):
    t = r.find(A + "t")
    txt = t.text if t is not None and t.text else ""
    if not txt:
        return ""
    rpr = r.find(A + "rPr")
    if rpr is not None:
        base = rpr.get("baseline")
        if base:
            try:
                b = int(base)
                if b > 0:   return f"<sup>{txt}</sup>"
                if b < 0:   return f"<sub>{txt}</sub>"
            except ValueError:
                pass
    return txt

def para_to_text(p):
    """Build one paragraph's text: runs + inline/display math in document order."""
    inline, display_only = [], True
    segs = []
    for node in p:
        tag = local(node)
        if tag == "pPr":
            continue
        if any(True for _ in node.iter(M + "oMath")):
            for lx in emit_math(node):
                if not lx:
                    continue
                ins, disp = wrap_math(lx)
                segs.append(("math_disp", disp) if disp else ("math_inl", ins))
        elif tag == "r":
            segs.append(("txt", run_text(node)))
            display_only = False
        elif tag == "br":
            segs.append(("txt", " "))
        elif tag == "fld":
            t = node.find(A + "t")
            if t is not None and t.text:
                segs.append(("txt", t.text)); display_only = False
        elif tag == "AlternateContent":
            # non-math AlternateContent: take Choice text
            choice = node.find(MC + "Choice")
            if choice is not None:
                for t in choice.iter(A + "t"):
                    if t.text:
                        segs.append(("txt", t.text)); display_only = False
    # if the paragraph is a single display-math block, return it raw
    txt_segs = [s for k, s in segs if k == "txt" and s.strip()]
    disp_segs = [s for k, s in segs if k == "math_disp"]
    if disp_segs and not txt_segs:
        return "\n".join(disp_segs)
    # otherwise assemble left-to-right, downgrading display math to inline
    out = []
    for k, s in segs:
        if k == "math_disp":
            out.append("$" + s.strip("$") + "$")
        else:
            out.append(s)
    return "".join(out).strip()

# ---------- shape geometry / meta ----------
def xfrm_of(sp):
    x = sp.find(f"{P}spPr/{A}xfrm") if sp.find(f"{P}spPr/{A}xfrm") is not None else None
    if x is None:
        # pictures / frames
        for path in (f"{P}spPr/{A}xfrm", f"{P}grpSpPr/{A}xfrm", f"{A}xfrm"):
            x = sp.find(path)
            if x is not None:
                break
    if x is None:
        return None
    off = x.find(A + "off"); ext = x.find(A + "ext")
    if off is None or ext is None:
        return None
    return (int(off.get("x", 0)), int(off.get("y", 0)),
            int(ext.get("cx", 0)), int(ext.get("cy", 0)))

def is_title_ph(sp):
    ph = sp.find(f"{P}nvSpPr/{P}nvPr/{P}ph")
    return ph is not None and ph.get("type") in ("title", "ctrTitle")

def max_font_pt(sp):
    m = None
    for rpr in sp.iter(A + "rPr"):
        sz = rpr.get("sz")
        if sz:
            try: m = max(m or 0, int(sz) / 100.0)
            except ValueError: pass
    return m

def has_bullets(txbody):
    for p in txbody.findall(A + "p"):
        ppr = p.find(A + "pPr")
        if ppr is None:
            continue
        lvl = ppr.get("lvl")
        if lvl and int(lvl) > 0:
            return True
        if ppr.find(A + "buChar") is not None or ppr.find(A + "buAutoNum") is not None:
            return True
    return False

# ---------- classification (same rules as v1) ----------
def on_image(tb, images):
    cx, cy = tb["x"] + tb["w"] / 2, tb["y"] + tb["h"] / 2
    for im in images:
        if im.get("file") and im["x"] - 1 <= cx <= im["x"] + im["w"] + 1 \
           and im["y"] - 1 <= cy <= im["y"] + im["h"] + 1:
            return True
    return False

def classify(text, tb, images, font):
    words = len(text.split())
    if ATTR.search(text) or (URLRE.search(text) and words <= 8) or AUTHYEAR.match(text):
        return "caption"
    if words <= 6 and on_image(tb, images):
        return "caption"
    if font is not None and font <= 11 and words <= 6:
        return "caption"
    return "body"

def guess_template(has_title, nbody, images):
    ni = len(images)
    if ni >= 1 and nbody == 0 and not has_title: return "T4-full-figure"
    if ni >= 1 and nbody == 0: return "T5-figure-focus"
    if ni == 2 and nbody <= 2: return "T6-two-image"
    if ni >= 1 and nbody > 0: return "T3-text-and-image"
    if ni == 0 and nbody > 0: return "T1-prose"
    return "T0-title-or-other"

# ---------- shape iteration (descends into mc:Choice) ----------
def iter_shapes(container):
    """Yield ('sp', el) / ('pic', el) in document order, taking Choice over Fallback."""
    for node in container:
        tag = local(node)
        if tag == "AlternateContent":
            choice = node.find(MC + "Choice")
            if choice is not None:
                yield from iter_shapes(choice)
        elif tag == "sp":
            yield ("sp", node)
        elif tag == "pic":
            yield ("pic", node)
        elif tag == "grpSp":
            yield from iter_shapes(node)
        elif tag == "graphicFrame":
            yield ("frame", node)

def extract(pptx_path, out_dir):
    pptx_path = pathlib.Path(pptx_path)
    name = pptx_path.stem
    out = pathlib.Path(out_dir) / name
    img_dir = out / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))
    W, H = prs.slide_width, prs.slide_height
    manifest = {"source": pptx_path.name, "aspect": round(W / H, 3),
                "n_slides": len(prs.slides), "slides": []}

    for idx, slide in enumerate(prs.slides, start=1):
        spTree = slide._element.find(f"{P}cSld/{P}spTree")
        texts, images, imgn = [], [], 0
        part = slide.part
        for kind, el in iter_shapes(spTree):
            geom = xfrm_of(el)
            gx, gy, gw, gh = geom if geom else (0, 0, 0, 0)
            if kind == "pic":
                blip = el.find(f"{P}blipFill/{A}blip")
                rid = blip.get(f"{{http://schemas.openxmlformats.org/officeDocument/2006/relationships}}embed") if blip is not None else None
                imgn += 1
                if rid:
                    try:
                        image_part = part.related_part(rid)
                        blob = image_part.blob
                        ext = image_part.partname.ext.lstrip(".")
                        ext, blob = optimize_image(blob, ext)   # downscale / convert TIFF / recompress
                        fn = f"slide{idx:03d}_img{imgn}.{ext}"
                        (img_dir / fn).write_bytes(blob)
                        images.append({"file": fn, "x": pct(gx, W), "y": pct(gy, H),
                                       "w": pct(gw, W), "h": pct(gh, H)})
                    except Exception as e:
                        images.append({"file": None, "error": str(e)})
                else:
                    images.append({"file": None, "error": "no embed"})
            elif kind == "sp":
                txbody = el.find(f"{P}txBody")
                if txbody is None:
                    continue
                lines = []
                for p in txbody.findall(A + "p"):
                    s = para_to_text(p)
                    if s.strip():
                        lines.append(s.strip())
                if not lines:
                    continue
                texts.append({"text": "\n".join(lines), "lines": lines,
                              "x": pct(gx, W), "y": pct(gy, H),
                              "w": pct(gw, W), "h": pct(gh, H),
                              "font": max_font_pt(el),
                              "_is_title": is_title_ph(el),
                              "bullets": has_bullets(txbody)})
        # assign roles now that images are known
        for tb in texts:
            tb["role"] = "title" if tb.pop("_is_title") else \
                classify(tb["text"].replace("\n", " "), tb, images, tb["font"])
        notes = ""
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
        nbody = sum(1 for t in texts if t["role"] == "body")
        has_title = any(t["role"] == "title" for t in texts)
        manifest["slides"].append({
            "n": idx, "source_ref": f"{pptx_path.name} slide {idx}",
            "texts": texts, "notes": notes, "images": images,
            "template_guess": guess_template(has_title, nbody, images),
        })
    (out / "manifest.json").write_text(json.dumps(manifest, indent=1))
    return manifest

if __name__ == "__main__":
    from collections import Counter
    m = extract(sys.argv[1], sys.argv[2])
    print(f"{m['source']}: {m['n_slides']} slides")
    print(" templates:", dict(Counter(s['template_guess'] for s in m['slides'])))
    roles = Counter(t['role'] for s in m['slides'] for t in s['texts'])
    print(" text roles:", dict(roles))
    nmath = sum(1 for s in m['slides'] for t in s['texts'] if "$" in t['text'])
    print(" text shapes containing math:", nmath)
    print(" slides with notes:", sum(1 for s in m['slides'] if s['notes']))
