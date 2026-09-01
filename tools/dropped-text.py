"""Find slides where the pptx converter dropped body text entirely.

The label-soup detector finds text that survived in the wrong place. This
finds text that did not survive at all. The converter classifies each slide
by template, and a slide read as "two images" can have its text boxes
discarded rather than scattered -- which leaves nothing behind to notice.
The Effect of Strain Rate slide in Module 3.3 lost three paragraphs that
way, and nothing in the markdown hinted at it.

The only way to see it is to go back to the source. Each deck carries the
converter's provenance comments (`pptx slide N`), so every slide can be
matched to its original and their word counts compared.

Reports a slide when the pptx has substantial text that the markdown does
not. Speaker notes count as a home for it -- text moved to the notes is not
lost. Titles, credits and numbers are ignored on both sides.

It compares VOCABULARY, not meaning, so a slide whose prose has been
deliberately rewritten reports as missing words even though nothing was
lost. Read the hits, do not act on the count: in Module 3 four of six were
real drops and two were rewrites.

Run:  pixi run python tools/dropped-text.py Lectures/Module-iii-lecture*-Theory.reveal.md
"""
import pathlib
import re
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("needs python-pptx (it is in the pixi env)")

PPTX_DIRS = ["migration/PPTs/Module3", "migration/PPTs/Module4"]
STOP = re.compile(r"^[\d\s.,;:()\[\]%°+\-–—/]*$")


def words(text):
    """Content words, so that punctuation and numbers do not pad the count."""
    return [w for w in re.findall(r"[A-Za-z][A-Za-z'’-]{2,}", text)]


def pptx_text(path):
    """-> {slide number: all its text} for one deck."""
    out = {}
    for i, slide in enumerate(Presentation(path).slides, 1):
        bits = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                bits.append(sh.text_frame.text)
            if sh.has_table:
                for row in sh.table.rows:
                    bits += [c.text for c in row.cells]
        out[i] = "\n".join(bits)
    return out


def deck_slides(path):
    """-> [(pptx number, markdown body incl. notes)] for one converted deck."""
    txt = pathlib.Path(path).read_text()
    txt = re.sub(r"^---\n.*?\n---\n", "", txt, flags=re.S)
    out = []
    for blk in re.split(r"\n<--o-->\n", txt):
        for vb in re.split(r"\n<--v-->\n", blk):
            m = re.search(r"pptx slide (\d+)", vb)
            if m:
                out.append((int(m.group(1)), vb))
    return out


def find_pptx(name):
    for d in PPTX_DIRS:
        p = pathlib.Path(d) / name
        if p.exists():
            return p
    return None


total = 0
for deck in sys.argv[1:]:
    slides = deck_slides(deck)
    src = re.search(r"pptx slide \d+", pathlib.Path(deck).read_text())
    name = re.search(r"source: (\S+\.pptx)", pathlib.Path(deck).read_text())
    if not name:
        print(f"{deck}: no provenance comments, skipping")
        continue
    ppt = find_pptx(name.group(1))
    if not ppt:
        print(f"{deck}: cannot find {name.group(1)}, skipping")
        continue
    original = pptx_text(ppt)
    print(f"\n=== {pathlib.Path(deck).name}  vs  {ppt.name}")
    for num, body in slides:
        if num not in original:
            continue
        # what the markdown has, notes included -- notes are a valid home
        have = set(w.lower() for w in words(re.sub(r"<[^>]+>|!\[[^\]]*\]\([^)]*\)",
                                                   " ", body)))
        want = words(original[num])
        missing = [w for w in want if w.lower() not in have]
        if len(missing) >= 25:
            m = re.search(r"^#{1,4} (.+)$", body, re.M)
            title = re.sub(r"&#?\w+;", "", m.group(1)).strip() if m else "(untitled)"
            print(f"  pptx s{num:<3} {len(missing):>3} words unaccounted for   {title[:44]}")
            print(f"            · {' '.join(missing[:18])} ...")
            total += 1
print(f"\n{total} slides may have lost text" if total else "\nnothing obviously dropped")
